import unittest
import json
from ParseCsv import parse_csv, convert_to_json
from ParseJson import parse_json
from ParsePptx import parse_pptx
from ParsePdf import format_data
from pptx import Presentation

class TestParseFunctions(unittest.TestCase):
    def setUp(self):
        # Create a sample CSV content
        self.csv_content = """id,name,role
E001,Alice,Personal Trainer
E002,Bob,Group Fitness Instructor
E003,Eve,Gym Manager
"""
        # Write the sample CSV content to a temporary file
        self.csv_file = 'test_dataset.csv'
        with open(self.csv_file, 'w', encoding='utf-8') as f:
            f.write(self.csv_content)

        # Create a sample JSON content
        self.json_content = {
            "companies": [
                {
                    "id": 1,
                    "name": "FitPro",
                    "industry": "Sports and Leisure",
                    "revenue": 95000000,
                    "location": "North America",
                    "employees": [
                        {"id": "E001", "name": "Alice", "role": "Personal Trainer", "cashmoneh": 45000, "hired_date": "2020-01-15"},
                        {"id": "E002", "name": "Bob", "role": "Group Fitness Instructor", "cashmoneh": 40000, "hired_date": "2019-06-22"},
                        {"id": "E003", "name": "Eve", "role": "Gym Manager", "cashmoneh": 70000, "hired_date": "2021-03-12"}
                    ],
                    "performance": {
                        "2023_Q1": {"revenue": 23000000, "profit_margin": 12.5},
                        "2023_Q2": {"revenue": 25000000, "profit_margin": 13.0}
                    }
                }
            ]
        }
        # Write the sample JSON content to a temporary file
        self.json_file = 'test_dataset.json'
        with open(self.json_file, 'w', encoding='utf-8') as f:
            json.dump(self.json_content, f, indent=2)

        # Create a sample PPTX content
        self.pptx_file = 'test_dataset.pptx'
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Sample Slide Title"
        content = slide.shapes.placeholders[1]
        content.text = "Sample Slide Content"
        presentation.save(self.pptx_file)

    def tearDown(self):
        # Remove the temporary files after tests
        import os
        os.remove(self.csv_file)
        os.remove(self.json_file)
        os.remove(self.pptx_file)

    def test_parse_csv(self):
        parsed_data = parse_csv(self.csv_file)
        expected_data = [
            {'id': 'E001', 'name': 'Alice', 'role': 'Personal Trainer'},
            {'id': 'E002', 'name': 'Bob', 'role': 'Group Fitness Instructor'},
            {'id': 'E003', 'name': 'Eve', 'role': 'Gym Manager'}
        ]
        self.assertEqual(parsed_data, expected_data)

    def test_convert_to_json(self):
        parsed_data = [
            {'id': 'E001', 'name': 'Alice', 'role': 'Personal Trainer'},
            {'id': 'E002', 'name': 'Bob', 'role': 'Group Fitness Instructor'},
            {'id': 'E003', 'name': 'Eve', 'role': 'Gym Manager'}
        ]
        json_data = convert_to_json(parsed_data)
        expected_json = json.dumps(parsed_data, indent=2)
        self.assertEqual(json_data, expected_json)

    def test_parse_json(self):
        parsed_data = parse_json(self.json_file)
        self.assertEqual(parsed_data, self.json_content)

    def test_format_data(self):
        raw_data = [
            "Year Quarter Revenue Memberships_Sold Avg_Duration",
            "2023 Q1 23,000,000 1500 12",
            "2023 Q2 25,000,000 1600 13"
        ]
        formatted_data = format_data(raw_data)
        expected_data = [
            {"year": 2023, "quarter": "Q1", "revenue": 23000000, "membershipsSold": 1500, "avgDuration": 12},
            {"year": 2023, "quarter": "Q2", "revenue": 25000000, "membershipsSold": 1600, "avgDuration": 13}
        ]
        self.assertEqual(formatted_data, expected_data)

    def test_parse_pptx(self):
        parsed_data = parse_pptx(self.pptx_file)
        expected_data = [
            {'title': 'Sample Slide Title', 'content': ['Sample Slide Content']}
        ]
        self.assertEqual(parsed_data, expected_data)

if __name__ == '__main__':
    unittest.main()