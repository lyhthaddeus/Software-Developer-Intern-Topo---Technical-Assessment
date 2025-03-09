from pptx import Presentation

def parse_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    
    parsed_data = []
    
    for slide in presentation.slides:
        slide_content = {'title': '', 'content': []}  # Dictionary to store slide content
        
        if slide.shapes.title:
            slide_content['title'] = slide.shapes.title.text.strip()  # Capture slide title
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():  # Check if shape contains text
                if shape != slide.shapes.title:  # Avoid capturing the title text again
                    slide_content['content'].append(shape.text.strip())  # Add text to the content list
            
            if hasattr(shape, 'table'):
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())  # Extract text from each cell
                    slide_content['content'].append(" | ".join(row_data))  # Join cell text with pipe symbol
        
        parsed_data.append(slide_content)
    
    return parsed_data

# Example usage
# pptx_file = "../datasets/dataset4.pptx"  # Replace with your PowerPoint file path
# parsed_data = parse_pptx(pptx_file)
#
# # Format the parsed data for frontend (as JSON)
# formatted_json = json.dumps(parsed_data, indent=2)
# print(formatted_json)
#

